"""MarkItUI Multi-Engine Worker.

Wird von der Electron-App aufgerufen und liefert seine Antwort als JSON auf stdout.
Damit Ausgaben fremder Bibliotheken (Warnungen von pdfminer, pptx, ...) die Antwort
nicht zerstoeren, wird das JSON in eindeutige Marker eingerahmt und waehrend der
Konvertierung alles andere auf stderr umgeleitet.
"""

import sys
import os
import io
import json
import argparse
import datetime
import re
import contextlib
from pathlib import Path

# UTF-8 erzwingen (Windows-Konsolen laufen sonst haeufig auf cp1252)
for _stream in ("stdout", "stderr"):
    try:
        getattr(sys, _stream).reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

JSON_BEGIN = "@@MARKITUI_JSON_BEGIN@@"
JSON_END = "@@MARKITUI_JSON_END@@"
# Markiert auf stderr, welche Engine gerade laeuft (fuer Absturz-Diagnose der App)
ENGINE_MARKER = "@@MARKITUI_ENGINE@@"

# Welche Pakete werden fuer welches Format gebraucht? (fuer verstaendliche Fehlermeldungen)
FORMAT_REQUIREMENTS = {
    ".pdf": ("PDF-Dateien", ["pdfplumber", "pypdfium2", "pdfminer.six"]),
    ".docx": ("Word-Dokumente", ["mammoth", "python-docx"]),
    ".doc": ("alte Word-Dokumente (.doc)", ["markitdown"]),
    ".pptx": ("PowerPoint-Praesentationen", ["python-pptx"]),
    ".ppt": ("alte PowerPoint-Dateien (.ppt)", ["markitdown"]),
    ".xlsx": ("Excel-Tabellen", ["openpyxl"]),
    ".xlsm": ("Excel-Tabellen mit Makros", ["openpyxl"]),
    ".xls": ("alte Excel-Tabellen (.xls)", ["xlrd"]),
    ".epub": ("E-Books", ["markitdown"]),
}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".log", ".ini", ".cfg", ".yaml", ".yml",
    ".json", ".xml", ".html", ".htm", ".csv", ".tsv", ".rtf", ".srt", ".vtt",
}

BINARY_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".xlsm",
    ".epub", ".odt", ".ods", ".odp", ".zip",
}


def emit(payload: dict) -> None:
    """Schreibt die Antwort eingerahmt auf stdout."""
    sys.stdout.write(JSON_BEGIN)
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))
    sys.stdout.write(JSON_END)
    sys.stdout.write("\n")
    sys.stdout.flush()


def format_title_from_stem(stem: str) -> str:
    clean = stem.replace("_", " ").replace("-", " ")
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean or stem


def yaml_quote(value: str) -> str:
    """Erzeugt einen gueltigen YAML-Doppelquote-String."""
    escaped = str(value).replace("\\", "\\\\").replace('"', '\\"')
    escaped = escaped.replace("\n", " ").replace("\r", " ")
    return f'"{escaped}"'


def module_available(module_name: str) -> bool:
    """Prueft, ob ein Modul importierbar ist, ohne es dauerhaft zu laden."""
    try:
        import importlib.util

        return importlib.util.find_spec(module_name) is not None
    except Exception:
        return False


# -------------------------------------------------------------
# Diagnose
# -------------------------------------------------------------

def check_environment() -> dict:
    """Diagnostiziert Python, Module und die Verfuegbarkeit der Konverter-Engines."""
    package_map = {
        "markitdown": "markitdown",
        "pdfplumber": "pdfplumber",
        "pypdfium2": "pypdfium2",
        "pdfminer": "pdfminer.six",
        "mammoth": "mammoth",
        "docx": "python-docx",
        "pptx": "python-pptx",
        "openpyxl": "openpyxl",
        "xlrd": "xlrd",
        "bs4": "beautifulsoup4",
        "markdown": "markdown",
        "pygments": "pygments",
    }

    installed = {}
    missing = []

    for module_name, package_name in package_map.items():
        ok = module_available(module_name)
        installed[package_name] = ok
        if not ok:
            missing.append(package_name)

    has_markitdown = installed.get("markitdown", False)
    has_pdf = any(installed.get(p, False) for p in ("pdfplumber", "pypdfium2", "pdfminer.six"))
    has_docx = installed.get("mammoth", False) or installed.get("python-docx", False)
    has_pptx = installed.get("python-pptx", False)
    has_xlsx = installed.get("openpyxl", False)

    version = sys.version_info
    python_ok = (version.major, version.minor) >= (3, 9)

    # "Bereit" heisst: alle Hauptformate koennen konvertiert werden - entweder ueber
    # MarkItDown oder ueber die spezialisierten Fallback-Engines.
    ready = python_ok and has_pdf and has_docx and has_pptx and has_xlsx

    return {
        "ready": ready,
        "python_version": ".".join(str(p) for p in version[:3]),
        "python_executable": sys.executable,
        "platform": sys.platform,
        "python_ok": python_ok,
        "installed_packages": [pkg for pkg, ok in installed.items() if ok],
        "missing_packages": missing,
        "has_markitdown": has_markitdown,
        "has_pdfplumber": has_pdf,
        "has_mammoth": installed.get("mammoth", False),
        "has_docx": has_docx,
        "has_pptx": has_pptx,
        "has_openpyxl": has_xlsx,
    }


# -------------------------------------------------------------
# Hilfsfunktionen
# -------------------------------------------------------------

def looks_binary(path: Path) -> bool:
    """Erkennt Binaerdateien anhand von Nullbytes / nicht darstellbaren Zeichen."""
    try:
        with open(path, "rb") as handle:
            chunk = handle.read(8192)
    except Exception:
        return False

    if not chunk:
        return False
    if b"\x00" in chunk:
        return True

    printable = sum(1 for byte in chunk if 32 <= byte < 127 or byte in (9, 10, 13))
    return (printable / len(chunk)) < 0.75


def read_text_file_safe(path: Path) -> str:
    """Liest eine Textdatei und probiert dabei gaengige Kodierungen durch."""
    raw_bytes = path.read_bytes()

    if raw_bytes.startswith(b"\xff\xfe") or raw_bytes.startswith(b"\xfe\xff"):
        try:
            return raw_bytes.decode("utf-16")
        except UnicodeDecodeError:
            pass

    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return raw_bytes.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue

    return raw_bytes.decode("utf-8", errors="replace")


def strip_rtf(text: str) -> str:
    """Entfernt RTF-Steuersequenzen und liefert lesbaren Text."""
    if "\\rtf" not in text[:200]:
        return text
    without_groups = re.sub(r"\{\\\*?[^{}]*\}", " ", text)
    without_controls = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", without_groups)
    cleaned = without_controls.replace("{", "").replace("}", "")
    return re.sub(r"[ \t]{2,}", " ", cleaned).strip()


_SENTENCE_END_RE = re.compile(r"[.!?:;…][\"'”’\)\]]*$")
_LIST_OR_HEADING_RE = re.compile(r"^\s*([-*+•‣▪◦]|\d+[.)]|\|.*\||#{1,6}\s|>\s)")
_LETTER_SPACED_WORD_RE = re.compile(r"(?:\b\w\s){3,}\w\b")


def _collapse_letter_spacing(line: str) -> str:
    """Manche PDF-Schriftarten liefern Buchstaben einzeln mit Leerzeichen dazwischen
    ("W o r t" statt "Wort"). Erkennt solche Laufweiten-Artefakte und zieht sie zusammen."""
    return _LETTER_SPACED_WORD_RE.sub(lambda m: m.group(0).replace(" ", ""), line)


def reflow_pdf_text(text: str) -> str:
    """Bereinigt typische PDF-Extraktionsartefakte bei Zeilen- und Wortabstaenden.

    PDF-Engines liefern Text so, wie er auf der Seite umgebrochen wurde - jede
    sichtbare Zeile endet mit einem Zeilenumbruch, unabhaengig davon, ob dort
    tatsaechlich ein Satz oder Absatz endet. Da die Vorschau harte Zeilenumbrueche
    als <br> rendert, wirkte importierter PDF-Text dadurch bisher wie eine Leiter
    aus lauter kurzen Zeilen mit unregelmaessigen Abstaenden. Hier werden weich
    umgebrochene Zeilen wieder zu Fliesstext zusammengefuehrt, Trennstriche am
    Zeilenende entfernt und doppelte/unregelmaessige Leerzeichen normalisiert.
    """
    if not text or not text.strip():
        return text

    raw_lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    # Erst Laufweiten-Artefakte erkennen (solange ein groesserer Abstand zwischen echten
    # Woertern noch von einfachen Buchstaben-Zwischenraeumen unterscheidbar ist), danach
    # erst alle mehrfachen Leer-/Tabzeichen auf ein einzelnes Leerzeichen normalisieren.
    lines = [_collapse_letter_spacing(ln.strip()) for ln in raw_lines]
    lines = [re.sub(r"[ \t]{2,}", " ", ln) for ln in lines]

    paragraphs = []
    buffer = ""

    def flush():
        nonlocal buffer
        if buffer:
            paragraphs.append(buffer)
            buffer = ""

    for line in lines:
        if not line:
            flush()
            if paragraphs and paragraphs[-1] != "":
                paragraphs.append("")
            continue

        if _LIST_OR_HEADING_RE.match(line):
            flush()
            paragraphs.append(line)
            continue

        if not buffer:
            buffer = line
            continue

        if buffer.endswith("-") and len(buffer) > 1 and buffer[-2].isalpha() and line[:1].islower():
            # Trennstrich, der nur wegen des Zeilenumbruchs eingefuegt wurde: Wort zusammenziehen.
            buffer = buffer[:-1] + line
        elif _SENTENCE_END_RE.search(buffer):
            # Zeile endet mit Satzzeichen -> vermutlich Absatzende, eigene Zeile behalten.
            paragraphs.append(buffer)
            buffer = line
        else:
            buffer += " " + line

    flush()

    result_lines = []
    for line in paragraphs:
        if line == "" and result_lines and result_lines[-1] == "":
            continue
        result_lines.append(line)

    return "\n".join(result_lines).strip()


NO_TEXT_LAYER_MARKER = "kein Text gefunden (evtl. ein reines Scan-PDF ohne Textebene)"


def rows_to_markdown_table(rows) -> list:
    """Formatiert Zeilenlisten als Markdown-Tabelle (mit Pipe-Escaping)."""
    lines = []
    if not rows:
        return lines

    def cell(value) -> str:
        text = "" if value is None else str(value)
        return text.replace("\n", " ").replace("|", "\\|").strip()

    header = [cell(c) for c in rows[0]]
    header = [h if h else f"Spalte {i + 1}" for i, h in enumerate(header)]

    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")

    for row in rows[1:]:
        cells = [cell(c) for c in row]
        while len(cells) < len(header):
            cells.append("")
        lines.append("| " + " | ".join(cells[: len(header)]) + " |")

    return lines


# -------------------------------------------------------------
# Spezialisierte Fallback-Konverter
# -------------------------------------------------------------

def convert_pdf_plumber(path: Path) -> str:
    """PDF ueber pdfplumber - liefert auch Tabellen als Markdown."""
    import pdfplumber

    pages_text = []
    with pdfplumber.open(str(path)) as pdf:
        page_count = len(pdf.pages)
        for idx, page in enumerate(pdf.pages, 1):
            try:
                tables = page.extract_tables()
            except Exception:
                tables = []
            text = reflow_pdf_text(page.extract_text() or "")

            page_content = []
            if page_count > 1:
                page_content.append(f"### Seite {idx}\n")
            if text.strip():
                page_content.append(text.strip())

            for table in tables or []:
                if not table or not any(table):
                    continue
                table_lines = rows_to_markdown_table(table)
                if table_lines:
                    page_content.append("\n" + "\n".join(table_lines) + "\n")

            if page_content:
                pages_text.append("\n\n".join(page_content))

    if not pages_text:
        raise RuntimeError(NO_TEXT_LAYER_MARKER)
    return "\n\n---\n\n".join(pages_text)


def convert_pdf_pdfium(path: Path) -> str:
    """PDF ueber pypdfium2 - schnell und ohne weitere Abhaengigkeiten."""
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(str(path))
    try:
        pages_text = []
        page_count = len(pdf)
        for idx in range(page_count):
            text = reflow_pdf_text(pdf[idx].get_textpage().get_text_range())
            if text and text.strip():
                if page_count > 1:
                    pages_text.append(f"### Seite {idx + 1}\n\n{text.strip()}")
                else:
                    pages_text.append(text.strip())
    finally:
        try:
            pdf.close()
        except Exception:
            pass

    if not pages_text:
        raise RuntimeError(NO_TEXT_LAYER_MARKER)
    return "\n\n---\n\n".join(pages_text)


def convert_pdf_pdfminer(path: Path) -> str:
    """PDF ueber pdfminer.six als letzte Reserve."""
    from pdfminer.high_level import extract_text

    text = reflow_pdf_text(extract_text(str(path)) or "")
    if not text or not text.strip():
        raise RuntimeError(NO_TEXT_LAYER_MARKER)
    return text.strip()


def convert_docx_mammoth(path: Path) -> str:
    """DOCX ueber mammoth - erhaelt Ueberschriften, Listen und Formatierung."""
    import mammoth

    with open(str(path), "rb") as docx_file:
        result = mammoth.convert_to_markdown(docx_file)

    if not result.value or not result.value.strip():
        raise RuntimeError("kein Inhalt gefunden")
    return result.value.strip()


def convert_docx_python_docx(path: Path) -> str:
    """DOCX ueber python-docx inklusive Tabellen."""
    import docx

    document = docx.Document(str(path))
    content = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name if paragraph.style else ""
        if style.startswith("Heading 1") or style.startswith("Ueberschrift 1"):
            content.append(f"# {text}")
        elif style.startswith("Heading 2") or style.startswith("Ueberschrift 2"):
            content.append(f"## {text}")
        elif style.startswith("Heading 3") or style.startswith("Ueberschrift 3"):
            content.append(f"### {text}")
        elif "List" in style or "Liste" in style:
            content.append(f"- {text}")
        else:
            content.append(text)

    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_lines = rows_to_markdown_table(rows)
        if table_lines:
            content.append("\n" + "\n".join(table_lines) + "\n")

    if not content:
        raise RuntimeError("kein Inhalt gefunden")
    return "\n\n".join(content)


def convert_pptx_fallback(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx ist nicht installiert.")

    try:
        presentation = Presentation(str(path))
    except Exception as err:
        raise RuntimeError(f"PPTX konnte nicht geoeffnet werden: {err}")

    slides_md = []

    for idx, slide in enumerate(presentation.slides, 1):
        slide_lines = []
        title_shape = None
        slide_title = f"Folie {idx}"

        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.text.strip():
                slide_title = title_shape.text.strip()
        except Exception:
            title_shape = None

        slide_lines.append(f"## Folie {idx}: {slide_title}\n")

        for shape in slide.shapes:
            if title_shape is not None and shape is title_shape:
                continue
            try:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            indent = "  " * (paragraph.level or 0)
                            slide_lines.append(f"{indent}- {text}")
                elif shape.has_table:
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    table_lines = rows_to_markdown_table(rows)
                    if table_lines:
                        slide_lines.append("")
                        slide_lines.extend(table_lines)
                        slide_lines.append("")
            except Exception:
                continue

        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"\n> **Notizen:** {notes}\n")
        except Exception:
            pass

        slides_md.append("\n".join(slide_lines).strip())

    if slides_md:
        return "\n\n---\n\n".join(slides_md)

    raise RuntimeError("Die Praesentation enthaelt keine lesbaren Folien.")


def convert_xlsx_fallback(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl ist nicht installiert.")

    workbook = None
    try:
        workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        sheets_md = []

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            rows = [
                row for row in worksheet.iter_rows(values_only=True)
                if row and any(cell is not None and str(cell).strip() != "" for cell in row)
            ]
            if not rows:
                continue

            sheet_lines = [f"### Tabelle: {sheet_name}\n"]
            sheet_lines.extend(rows_to_markdown_table(rows))
            sheets_md.append("\n".join(sheet_lines))

        if sheets_md:
            return "\n\n---\n\n".join(sheets_md)
    except Exception as err:
        raise RuntimeError(f"Excel-Konvertierung fehlgeschlagen: {err}")
    finally:
        if workbook is not None:
            try:
                workbook.close()
            except Exception:
                pass

    raise RuntimeError("Die Excel-Datei enthaelt keine lesbaren Daten.")


def convert_xls_fallback(path: Path) -> str:
    """Alte .xls-Dateien ueber xlrd."""
    try:
        import xlrd
    except ImportError:
        raise RuntimeError("xlrd ist nicht installiert (wird fuer alte .xls-Dateien benoetigt).")

    try:
        book = xlrd.open_workbook(str(path))
    except Exception as err:
        raise RuntimeError(f"XLS konnte nicht geoeffnet werden: {err}")

    sheets_md = []
    for sheet in book.sheets():
        rows = [sheet.row_values(r) for r in range(sheet.nrows)]
        rows = [row for row in rows if any(str(c).strip() for c in row)]
        if not rows:
            continue
        sheet_lines = [f"### Tabelle: {sheet.name}\n"]
        sheet_lines.extend(rows_to_markdown_table(rows))
        sheets_md.append("\n".join(sheet_lines))

    if sheets_md:
        return "\n\n---\n\n".join(sheets_md)

    raise RuntimeError("Die XLS-Datei enthaelt keine lesbaren Daten.")


def convert_csv_fallback(path: Path) -> str:
    import csv

    text = read_text_file_safe(path)
    lines = text.strip().splitlines()
    if not lines:
        return ""

    sample = "\n".join(lines[:20])
    try:
        delimiter = csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
    except Exception:
        delimiter = ";" if sample.count(";") > sample.count(",") else ","

    rows = list(csv.reader(lines, delimiter=delimiter))
    if not rows:
        return text

    return "\n".join(rows_to_markdown_table(rows))


def convert_html_fallback(path: Path) -> str:
    text = read_text_file_safe(path)
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        extracted = soup.get_text("\n")
        cleaned = re.sub(r"\n{3,}", "\n\n", extracted).strip()
        if cleaned:
            return cleaned
    except Exception:
        pass
    return text


# -------------------------------------------------------------
# Orchestrierung
# -------------------------------------------------------------

def convert_with_markitdown(path: Path) -> str:
    """Tier-1-Engine von Microsoft (deckt auch exotische Formate ab)."""
    from markitdown import MarkItDown

    result = MarkItDown().convert(str(path))
    text = result.text_content if hasattr(result, "text_content") else str(result)
    if not text or not text.strip():
        raise RuntimeError("MarkItDown hat keinen Text geliefert.")
    return text


def convert_pdf_via_markitdown(path: Path) -> str:
    """MarkItDown als PDF-Fallback - nutzt intern ebenfalls pdfminer und hat damit
    dieselben Zeilenumbruch-/Wortabstand-Artefakte wie die anderen PDF-Engines."""
    return reflow_pdf_text(convert_with_markitdown(path))


def convert_text(path: Path) -> str:
    return read_text_file_safe(path)


def convert_rtf(path: Path) -> str:
    return strip_rtf(read_text_file_safe(path))


def convert_image_placeholder(path: Path) -> str:
    return f"![{path.name}]({path.name})\n\n*Bilddokument: {path.name}*"


def engines_for(ext: str):
    """Engine-Kette fuer einen Dateityp - spezialisierte Engines zuerst.

    Die nativen Engines liefern bessere Tabellen/Notizen und starten deutlich
    schneller als MarkItDown (das beim Import ein Machine-Learning-Modell laedt).
    MarkItDown bleibt als universeller Auffang-Konverter am Ende der Kette.
    """
    if ext == ".pdf":
        return [
            ("pdfplumber", convert_pdf_plumber),
            ("pypdfium2", convert_pdf_pdfium),
            ("markitdown", convert_pdf_via_markitdown),
            ("pdfminer", convert_pdf_pdfminer),
        ]
    if ext == ".docx":
        return [
            ("mammoth", convert_docx_mammoth),
            ("python-docx", convert_docx_python_docx),
            ("markitdown", convert_with_markitdown),
        ]
    if ext == ".pptx":
        return [("python-pptx", convert_pptx_fallback), ("markitdown", convert_with_markitdown)]
    if ext in (".xlsx", ".xlsm"):
        return [("openpyxl", convert_xlsx_fallback), ("markitdown", convert_with_markitdown)]
    if ext == ".xls":
        return [("xlrd", convert_xls_fallback), ("markitdown", convert_with_markitdown)]
    if ext in (".csv", ".tsv"):
        return [("csv-parser", convert_csv_fallback), ("markitdown", convert_with_markitdown)]
    if ext in (".html", ".htm"):
        return [("html-parser", convert_html_fallback), ("markitdown", convert_with_markitdown)]
    if ext == ".rtf":
        return [("rtf-decoder", convert_rtf), ("markitdown", convert_with_markitdown)]
    if ext in TEXT_EXTENSIONS:
        return [("text-decoder", convert_text), ("markitdown", convert_with_markitdown)]
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"):
        return [("markitdown", convert_with_markitdown), ("image-embed", convert_image_placeholder)]

    # Unbekannt (.doc, .ppt, .epub, ...): MarkItDown kann am meisten abdecken.
    chain = [("markitdown", convert_with_markitdown)]
    if ext == ".doc":
        chain.append(("rtf-decoder", convert_rtf))
    return chain


def skipped_engines() -> set:
    """Engines, die in einem vorherigen Versuch den Interpreter zum Absturz gebracht haben."""
    raw = os.environ.get("MARKITUI_SKIP_ENGINES", "")
    return {name.strip().lower() for name in raw.split(",") if name.strip()}


def missing_packages_for(ext: str) -> list:
    """Liefert die fuer dieses Format fehlenden Pakete."""
    entry = FORMAT_REQUIREMENTS.get(ext)
    if not entry:
        return []

    module_for_package = {
        "pdfplumber": "pdfplumber",
        "pypdfium2": "pypdfium2",
        "pdfminer.six": "pdfminer",
        "mammoth": "mammoth",
        "python-docx": "docx",
        "python-pptx": "pptx",
        "openpyxl": "openpyxl",
        "xlrd": "xlrd",
        "markitdown": "markitdown",
    }

    return [
        package for package in entry[1]
        if not module_available(module_for_package.get(package, package))
    ]


def convert_document(file_path: str, add_frontmatter: bool = True, tags=None,
                     subject: str = "", title: str = "") -> dict:
    path = Path(file_path)
    if not path.exists():
        return {
            "success": False,
            "markdown": "",
            "error": f"Datei existiert nicht: {file_path}",
            "file_name": path.name,
        }

    ext = path.suffix.lower()
    raw_text = ""
    used_engine = ""
    problems = []
    skip = skipped_engines()
    # PDFs ohne Textebene (reine Scans) schlagen bei jeder Extraktions-Engine mit
    # derselben Ursache fehl - keine der Engines kann OCR. Sobald zwei Engines das
    # uebereinstimmend melden, lohnt es sich nicht mehr, auch noch MarkItDown (laedt
    # ein ML-Modell) und pdfminer durchzuprobieren: das kostet nur Zeit/Speicher und
    # aendert am Ergebnis nichts. Stattdessen wird sofort mit einer klaren Diagnose
    # abgebrochen.
    no_text_layer_hits = 0

    for name, engine in engines_for(ext):
        if name.lower() in skip:
            problems.append(f"{name}: uebersprungen (vorheriger Absturz)")
            continue

        # Marker fuer die App: stuerzt der Interpreter ab (z. B. defekte native
        # Bibliothek), weiss sie, welche Engine schuld war, und wiederholt ohne sie.
        sys.stderr.write(f"{ENGINE_MARKER}{name}\n")
        sys.stderr.flush()

        try:
            text = engine(path)
        except ImportError as err:
            problems.append(f"{name}: nicht installiert ({err})")
            continue
        except Exception as err:
            problems.append(f"{name}: {err}")
            if ext == ".pdf" and str(err) == NO_TEXT_LAYER_MARKER:
                no_text_layer_hits += 1
                if no_text_layer_hits >= 2:
                    problems.append(
                        "Weitere Engines uebersprungen: uebereinstimmend keine Textebene gefunden "
                        "(vermutlich ein gescanntes PDF, das OCR benoetigt)."
                    )
                    break
            continue

        if text and text.strip():
            raw_text = text
            used_engine = name
            break
        problems.append(f"{name}: kein Inhalt gefunden")

    # Letzter Ausweg fuer Textformate ohne Ergebnis
    if not raw_text.strip() and ext not in BINARY_EXTENSIONS and not looks_binary(path):
        try:
            raw_text = read_text_file_safe(path)
            used_engine = "text-decoder"
        except Exception as err:
            problems.append(f"text-decoder: {err}")

    if not raw_text or not raw_text.strip():
        missing = missing_packages_for(ext)
        label = FORMAT_REQUIREMENTS.get(ext, ("diesen Dateityp", []))[0]
        if ext == ".pdf" and no_text_layer_hits >= 2 and not missing:
            error_msg = (
                f"„{path.name}“ enthaelt keine Textebene - vermutlich ein gescanntes PDF "
                "(nur Bilder der Seiten, kein eingebetteter Text).\n\n"
                "MarkItUI kann daraus derzeit keinen Text extrahieren (keine OCR-Engine installiert). "
                "Bitte das PDF vorher mit einem OCR-Werkzeug (z. B. Adobe Acrobat, OCRmyPDF) durchsuchbar machen."
            )
        elif missing:
            error_msg = (
                f"Fuer {label} fehlen Python-Pakete: {', '.join(missing)}.\n\n"
                'Klicke auf "1-Klick Pakete reparieren / installieren", dann klappt die Umwandlung.\n\n'
                "Details:\n" + "\n".join(problems)
            )
        else:
            error_msg = (
                f"Das Dokument konnte nicht umgewandelt werden ({path.name}).\n\n"
                + "\n".join(problems)
            )
        return {
            "success": False,
            "markdown": "",
            "error": error_msg,
            "file_name": path.name,
            "missing_prerequisites": bool(missing),
        }

    # Sicherheitsnetz: niemals Binaermuell als "erfolgreich" zurueckgeben
    if ext in BINARY_EXTENSIONS and used_engine == "text-decoder":
        missing = missing_packages_for(ext)
        label = FORMAT_REQUIREMENTS.get(ext, ("diesen Dateityp", []))[0]
        return {
            "success": False,
            "markdown": "",
            "error": (
                f"Fuer {label} fehlt eine passende Engine"
                + (f" ({', '.join(missing)})" if missing else "")
                + '.\nKlicke auf "1-Klick Pakete reparieren / installieren".'
            ),
            "file_name": path.name,
            "missing_prerequisites": True,
        }

    # Aufbereitung
    cleaned_text = raw_text.strip()
    doc_title = title.strip() if title and title.strip() else format_title_from_stem(path.stem)

    lines = cleaned_text.splitlines()
    clean_lines = []
    found_first_heading = False

    for line in lines:
        stripped = line.strip()
        if not found_first_heading and stripped:
            match = re.match(r"^(#{1,3}\s*)?[Tt]itle:\s*(.+)$", stripped)
            if match:
                clean_lines.append(f"# {match.group(2).strip()}")
                found_first_heading = True
                continue
            if stripped.startswith("#"):
                found_first_heading = True
        clean_lines.append(line)

    cleaned_text = "\n".join(clean_lines).strip()

    if not cleaned_text.startswith("#"):
        cleaned_text = f"# {doc_title}\n\n{cleaned_text}"

    if add_frontmatter:
        now = datetime.datetime.now()
        tag_list = [str(t) for t in (tags or ["schule", "itslearning"])]
        if subject and subject.strip() and subject.strip().lower() not in [t.lower() for t in tag_list]:
            tag_list.append(subject.strip().lower())

        yaml_lines = [
            "---",
            f"title: {yaml_quote(doc_title)}",
            f"date: {now.strftime('%Y-%m-%d')}",
            f"created: {now.strftime('%Y-%m-%d %H:%M')}",
            f"source_file: {yaml_quote(path.name)}",
            f"source_type: {yaml_quote(path.suffix.lstrip('.'))}",
        ]
        if subject:
            yaml_lines.append(f"subject: {yaml_quote(subject)}")

        cleaned_tags = [t.strip().lstrip("#") for t in tag_list if t and t.strip().lstrip("#")]
        if cleaned_tags:
            yaml_lines.append("tags:")
            yaml_lines.extend(f"  - {t}" for t in cleaned_tags)

        yaml_lines.append("---")
        final_text = "\n".join(yaml_lines) + "\n\n" + cleaned_text + "\n"
    else:
        final_text = cleaned_text + "\n"

    return {
        "success": True,
        "markdown": final_text,
        "error": None,
        "file_name": path.name,
        "char_count": len(final_text),
        "engine_used": used_engine,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="MarkItUI Multi-Engine Worker")
    parser.add_argument("--file", type=str, help="Pfad zur Quelldatei")
    parser.add_argument("--frontmatter", action="store_true", default=True)
    parser.add_argument("--no-frontmatter", action="store_false", dest="frontmatter")
    parser.add_argument("--tags", type=str, default="schule,itslearning")
    parser.add_argument("--subject", type=str, default="")
    parser.add_argument("--title", type=str, default="")
    parser.add_argument("--json-input", action="store_true", help="Eingabe als JSON ueber stdin")
    parser.add_argument("--doctor", action="store_true", help="Diagnose als JSON ausgeben")

    args = parser.parse_args()

    if args.doctor:
        emit(check_environment())
        return

    if args.json_input:
        try:
            raw_input_data = sys.stdin.read()
        except Exception as err:
            emit({"success": False, "markdown": "", "error": f"stdin nicht lesbar: {err}", "file_name": ""})
            return

        if not raw_input_data.strip():
            emit({"success": False, "markdown": "", "error": "Leere Eingabe ueber stdin", "file_name": ""})
            return

        try:
            input_data = json.loads(raw_input_data)
        except Exception as err:
            emit({"success": False, "markdown": "", "error": f"Ungueltige JSON-Eingabe: {err}", "file_name": ""})
            return

        file_path = input_data.get("file_path", "")
        try:
            # Alles, was Bibliotheken auf stdout schreiben, landet auf stderr -
            # sonst zerstoert es die JSON-Antwort.
            buffer = io.StringIO()
            with contextlib.redirect_stdout(buffer):
                result = convert_document(
                    file_path,
                    input_data.get("add_frontmatter", True),
                    input_data.get("tags", ["schule", "itslearning"]),
                    input_data.get("subject", ""),
                    input_data.get("title", ""),
                )
            noise = buffer.getvalue()
            if noise.strip():
                sys.stderr.write(noise)
            emit(result)
        except MemoryError:
            emit({
                "success": False,
                "markdown": "",
                "error": "Die Datei ist zu gross fuer den verfuegbaren Arbeitsspeicher.",
                "file_name": os.path.basename(file_path),
            })
        except Exception as err:
            emit({
                "success": False,
                "markdown": "",
                "error": f"Unerwarteter Fehler: {err}",
                "file_name": os.path.basename(file_path),
            })
        return

    if args.file:
        tags = [t.strip() for t in args.tags.split(",") if t.strip()]
        buffer = io.StringIO()
        with contextlib.redirect_stdout(buffer):
            result = convert_document(args.file, args.frontmatter, tags, args.subject, args.title)
        emit(result)
        return

    parser.print_help()


if __name__ == "__main__":
    main()
